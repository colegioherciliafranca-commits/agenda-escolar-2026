#!/usr/bin/env python3
import os
import json
import uuid
from flask import Flask, request, jsonify, send_from_directory, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
import shutil
from pptx import Presentation
import re
from PIL import Image
import io
from vlibras_service import translate_text_for_libras, get_vlibras_service
import threading
import time

app = Flask(__name__, static_folder='client/build', static_url_path='')
CORS(app)

# Configuração
UPLOAD_FOLDER = 'uploads'
PROCESSED_FOLDER = 'processed'
ALLOWED_EXTENSIONS = {'ppt', 'pptx'}
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['PROCESSED_FOLDER'] = PROCESSED_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

# Criar diretórios necessários
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PROCESSED_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def simplify_text(text):
    """Simplifica o texto para torná-lo mais acessível"""
    text = re.sub(r'\s+', ' ', text).strip()
    
    sentences = re.split(r'[.!?]+', text)
    simplified_sentences = []
    
    for sentence in sentences:
        sentence = sentence.strip()
        if len(sentence) > 100:
            words = sentence.split()
            mid = len(words) // 2
            simplified_sentences.append(' '.join(words[:mid]) + '.')
            simplified_sentences.append(' '.join(words[mid:]) + '.')
        elif sentence:
            simplified_sentences.append(sentence + '.')
    
    return ' '.join(simplified_sentences)

def extract_slide_content(slide, slide_number, session_id):
    """Extrai conteúdo de um slide"""
    content = {
        'title': '',
        'text': '',
        'simplified_text': '',
        'images': [],
        'notes': '',
        'slide_number': slide_number
    }
    
    # Criar diretório para imagens do slide
    slide_images_dir = os.path.join(PROCESSED_FOLDER, session_id, f'slide_{slide_number}')
    os.makedirs(slide_images_dir, exist_ok=True)
    
    # Extrair título e texto
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text and not content['title']:
                content['title'] = text
            elif text and text != content['title']:
                content['text'] += text + ' '
    
    # Extrair notas
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text
        content['notes'] = notes_text.strip()
    
    # Simplificar texto
    full_text = content['title'] + ' ' + content['text']
    content['simplified_text'] = simplify_text(full_text)
    
    # Extrair e salvar imagens
    image_count = 0
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            try:
                image_data = shape.image.blob
                img = Image.open(io.BytesIO(image_data))
                
                # Salvar imagem
                image_filename = f'image_{image_count + 1}.png'
                image_path = os.path.join(slide_images_dir, image_filename)
                
                # Converter para RGB se necessário e salvar
                if img.mode in ('RGBA', 'LA', 'P'):
                    rgb_img = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    rgb_img.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                    rgb_img.save(image_path, 'PNG')
                else:
                    img.save(image_path, 'PNG')
                
                content['images'].append({
                    'filename': image_filename,
                    'width': img.width,
                    'height': img.height,
                    'size': len(image_data),
                    'url': f'/api/images/{session_id}/slide_{slide_number}/{image_filename}'
                })
                
                image_count += 1
            except Exception as e:
                print(f"Erro ao processar imagem: {e}")
                pass
    
    return content

def process_powerpoint(file_path, session_id):
    """Processa arquivo PowerPoint"""
    try:
        output_dir = os.path.join(PROCESSED_FOLDER, session_id)
        os.makedirs(output_dir, exist_ok=True)
        
        prs = Presentation(file_path)
        slides_data = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = extract_slide_content(slide, i + 1, session_id)
            slides_data.append(slide_content)
        
        # Salvar dados dos slides
        output_file = os.path.join(output_dir, 'slides.json')
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump({
                'total_slides': len(slides_data),
                'slides': slides_data
            }, f, ensure_ascii=False, indent=2)
        
        return {
            'success': True,
            'total_slides': len(slides_data),
            'slides': slides_data
        }
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e)
        }

@app.route('/')
def serve_index():
    return send_from_directory(app.static_folder, 'index.html')

@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'presentation' not in request.files:
        return jsonify({'error': 'Nenhum arquivo enviado'}), 400
    
    file = request.files['presentation']
    if file.filename == '':
        return jsonify({'error': 'Nenhum arquivo selecionado'}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        session_id = str(uuid.uuid4())
        result = process_powerpoint(file_path, session_id)
        
        if result['success']:
            return jsonify({
                'success': True,
                'sessionId': session_id,
                'slides': result['slides'],
                'message': 'Apresentação processada com sucesso'
            })
        else:
            return jsonify({'error': result['error']}), 500
    
    return jsonify({'error': 'Tipo de arquivo não permitido'}), 400

@app.route('/api/slides/<session_id>')
def get_slides(session_id):
    slides_path = os.path.join(PROCESSED_FOLDER, session_id, 'slides.json')
    
    if os.path.exists(slides_path):
        with open(slides_path, 'r', encoding='utf-8') as f:
            slides_data = json.load(f)
        return jsonify({'slides': slides_data['slides']})
    else:
        return jsonify({'error': 'Sessão não encontrada'}), 404

@app.route('/api/translate/<session_id>/<slide_number>', methods=['POST'])
def translate_slide(session_id, slide_number):
    """Traduz o texto de um slide específico para Libras em tempo real"""
    try:
        # Obter dados do slide
        slides_path = os.path.join(PROCESSED_FOLDER, session_id, 'slides.json')
        if not os.path.exists(slides_path):
            return jsonify({'error': 'Sessão não encontrada'}), 404
        
        with open(slides_path, 'r', encoding='utf-8') as f:
            slides_data = json.load(f)
        
        # Encontrar o slide específico
        slide_num = int(slide_number)
        slide = None
        for s in slides_data['slides']:
            if s['slide_number'] == slide_num:
                slide = s
                break
        
        if not slide:
            return jsonify({'error': 'Slide não encontrado'}), 404
        
        # Usar texto simplificado para tradução
        text_to_translate = slide.get('simplified_text', slide.get('text', ''))
        if not text_to_translate:
            return jsonify({'error': 'Sem texto para traduzir'}), 400
        
        # Verificar se já existe tradução em cache
        translation_cache_path = os.path.join(PROCESSED_FOLDER, session_id, f'translations_{slide_num}.json')
        if os.path.exists(translation_cache_path):
            with open(translation_cache_path, 'r', encoding='utf-8') as f:
                cached_data = json.load(f)
                # Se a tradução tiver menos de 1 hora, usar do cache
                if time.time() - cached_data.get('timestamp', 0) < 3600:
                    return jsonify({
                        'success': True,
                        'method': cached_data.get('method', 'widget'),
                        'widget_ready': cached_data.get('widget_ready', False),
                        'text_translated': cached_data.get('text_translated'),
                        'from_cache': True
                    })
        
        # Iniciar tradução em background
        def translate_in_background():
            try:
                translation_result = translate_text_for_libras(text_to_translate)
                
                # Salvar em cache com resultado do Widget
                cache_data = {
                    'method': translation_result.get('method', 'widget'),
                    'text_translated': text_to_translate,
                    'widget_ready': True,
                    'timestamp': time.time(),
                    'slide_number': slide_num,
                    'success': True
                }
                
                with open(translation_cache_path, 'w', encoding='utf-8') as f:
                    json.dump(cache_data, f, ensure_ascii=False, indent=2)
                    
            except Exception as e:
                print(f"Erro na tradução em background: {e}")
        
        # Iniciar thread para tradução assíncrona
        thread = threading.Thread(target=translate_in_background)
        thread.daemon = True
        thread.start()
        
        return jsonify({
            'success': True,
            'message': 'Tradução iniciada em background',
            'processing': True
        })
        
    except Exception as e:
        return jsonify({'error': f'Erro na tradução: {str(e)}'}), 500

@app.route('/api/translation-status/<session_id>/<slide_number>')
def get_translation_status(session_id, slide_number):
    """Verifica o status da tradução de um slide"""
    try:
        translation_cache_path = os.path.join(PROCESSED_FOLDER, session_id, f'translations_{slide_number}.json')
        
        if os.path.exists(translation_cache_path):
            with open(translation_cache_path, 'r', encoding='utf-8') as f:
                cached_data = json.load(f)
                
            if cached_data.get('widget_ready') or cached_data.get('video_url'):
                return jsonify({
                    'success': True,
                    'completed': True,
                    'method': cached_data.get('method', 'widget'),
                    'widget_ready': cached_data.get('widget_ready', False),
                    'video_url': cached_data.get('video_url'),
                    'text_translated': cached_data.get('text_translated')
                })
            else:
                return jsonify({
                    'success': True,
                    'completed': False,
                    'processing': True
                })
        else:
            return jsonify({
                'success': True,
                'completed': False,
                'not_started': True
            })
            
    except Exception as e:
        return jsonify({'error': f'Erro ao verificar status: {str(e)}'}), 500

@app.route('/api/images/<session_id>/<slide_folder>/<filename>')
def serve_image(session_id, slide_folder, filename):
    image_path = os.path.join(PROCESSED_FOLDER, session_id, slide_folder, filename)
    if os.path.exists(image_path):
        return send_file(image_path, mimetype='image/png')
    else:
        return jsonify({'error': 'Imagem não encontrada'}), 404

@app.route('/avatar')
def serve_avatar():
    return send_file('avatar_completo.html')

@app.route('/<path:path>')
def serve_static(path):
    return send_from_directory(app.static_folder, path)

if __name__ == '__main__':
    print("Servidor iniciado em http://localhost:5000")
    print("Acesse a aplicação no navegador!")
    app.run(debug=True, host='0.0.0.0', port=5000)
